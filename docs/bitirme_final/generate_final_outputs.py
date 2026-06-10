from __future__ import annotations

import re
import zipfile
from pathlib import Path

from docx import Document
from docx.enum.section import WD_SECTION_START
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT, WD_CELL_VERTICAL_ALIGNMENT
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches as PptInches, Pt as PptPt


ROOT = Path(__file__).resolve().parents[2]
FINAL_DIR = ROOT / "docs" / "bitirme_final"
SOURCE_DIR = ROOT / "docs" / "yazilim_muh_final"

MASTER_MD = FINAL_DIR / "BITIRME_FINAL_MASTER.md"
POSTER_MD = FINAL_DIR / "BITIRME_FINAL_POSTER.md"
REPORT_DOCX = FINAL_DIR / "Bitirme_Final_Raporu.docx"
POSTER_PPTX = FINAL_DIR / "Bitirme_Final_Posteri.pptx"
NOTES_MD = FINAL_DIR / "TESLIM_NOTLARI.md"


FINAL_SOURCE_ORDER = [
    "01_ANA_RAPOR.md",
    "02_KULLANIM_SENARYOLARI_VE_SOZLESMELER.md",
    "04_GEREKSINIM_IZLENEBILIRLIK_VE_VERI_SOZLUGU.md",
    "09_MODUL_BAZLI_DETAYLI_ANALIZ.md",
    "13_API_SOZLESME_VE_ENDPOINT_ANALIZI.md",
    "14_DETAYLI_VERITABANI_TASARIMI.md",
    "10_EKRAN_TASARIMLARI_VE_ARAYUZ_AKISLARI.md",
    "11_DETAYLI_TEST_SENARYOLARI.md",
    "03_TEST_RISK_PLAN_EKLERI.md",
    "12_GENISLETILMIS_RISK_YONETIMI.md",
    "06_MALIYET_KESTIRIMI_VE_KAYNAK_PLANLAMA.md",
    "07_KURULUM_EGITIM_BAKIM_VE_OPERASYON_PLANI.md",
    "08_GRAFIKLER_VE_DIYAGRAMLAR_KULLANIM_REHBERI.md",
]


RECENT_COMPLETION = """\
## Guncel Tamamlama Notu - 09.06.2026

Bu bolum, rapor paketinin ilk hazirlanmasindan sonra tamamlanan ve bitirme final raporuna eklenmesi gereken son proje gelistirmelerini ozetler.

### Veteriner Modulu Son Durum

- Veteriner arama ekranina tur ve hizmet filtresi eklendi.
- Backend listeleme sorgulari `species` ve `service` parametrelerini destekleyecek sekilde guncellendi.
- Randevu detay ekranina yeniden planlama arayuzu eklendi.
- Randevu yeniden planlama akisi backend `PATCH /api/appointments/:id/reschedule` endpoint'i ile baglandi.
- Online randevular icin sahte `meet.google.com/...` linki otomatik uretimi kaldirildi.
- Veteriner, online randevuyu onaylarken gercek gorusme linki girebilecek hale getirildi.
- Giris yapmadan veteriner mesajlasmasina gidildiginde, login sonrasi ilgili veteriner detayina geri donus akisi duzeltildi.

### Bakici Modulu Son Durum

- Bakici profil fotografi yuklemek icin backend `POST /api/pet-sitters/:id/avatar` endpoint'i eklendi.
- Flutter bakici profil olusturma/duzenleme akisi avatar yukleme endpoint'ine baglandi.
- Bakici finansal ozet ekrani dashboard uzerinden erisilebilir hale getirildi.
- Bakici rezervasyon ekraninda calisma saatleri gosterildi.
- Calisma saatleri disinda rezervasyon olusturma mobil arayuzde engellendi.
- Bakici yorumlarinda ayni kullanicinin ayni bakiciya birden fazla yorum yazmasi backend tarafinda engellendi.
- Bakim raporu fotograf yukleme ve canli konum takip akislarinin backend ve mobil baglantilari kontrol edildi.

### Dogrulama ve Kurulum Sonuclari

- Backend Jest testleri calistirildi: 2 test suite ve 10 test basarili.
- Flutter widget testleri calistirildi: bakici finansal ozet, bakim raporu detayi ve veteriner kazanc ekranlari basarili.
- Android debug APK basariyla derlendi.
- Dogru mobil proje uzerinden uretilen APK fiziksel Android cihaza ADB ile yeniden kuruldu.
- `flutter analyze` komutu Flutter analyzer internal crash verdigi icin sonuc uretmedi; buna karsin build ve testler basarili tamamlandi.

### Kalan Not

Gercek Google Meet linkinin otomatik uretilmesi icin Google Calendar/Meet OAuth entegrasyonu gerekir. Mevcut uygulamada sahte link uretilmemekte, veteriner tarafindan girilen gercek link saklanmaktadir.
"""


OLD_REPORT_SUMMARY = """\
## Onceki Rapor Kaynaklari ve Birlesim Notu

Projede daha once hazirlanan `RAPOR.md`, `BITIRME_RAPORU.md`, `BITIRME_RAPORU.txt` ve `docs/bitirme2_ara_rapor_guncel.md` dosyalari incelenmistir. Bu belgeler proje gecmisi, ara rapor anlatimi, okul form sablonu ve ekran goruntusu yer tutuculari acisindan kaynak olarak kullanilabilir. Ancak final raporda tekrar olusmamasi icin ana govde `docs/yazilim_muh_final/` klasorundeki kapsamli final paketinden uretilmistir.

Eski raporlarin final belgeye kattigi ana bilgiler sunlardir:

- Bitirme-1 ve Bitirme-2 surecinde proje kapsamının nasil genisledigi.
- Canli sunucu, MongoDB Atlas, mobil uygulama, admin paneli ve satici paneli gibi katmanlarin proje boyunca olgunlasmasi.
- Veteriner, bakici, magaza, sosyal akis, kayip/bulunan, sahiplendirme, mesajlasma ve bildirim modullerinin donemsel gelisim ozeti.
- Rapor icin kullanilabilecek ekran goruntusu basliklari.

Bu nedenle eski raporlar ek kaynak olarak saklanmis, final raporun ana akisi ise tek ve tutarli bir akademik rapor duzeninde yeniden toparlanmistir.
"""


FIGURE_INDEX = """\
## Sekil ve Diyagram Kaynaklari

Final rapor paketinde Word'e veya postere eklenebilecek hazir gorsel kaynaklar bulunmaktadir.

### Diyagram SVG Ciktilari

- `docs/yazilim_muh_final/rendered/diagrams/`: 25 adet UML ve sistem diyagrami.
- Oncelikli diyagramlar: use case, context, data flow, deployment, component, class, MongoDB view, veteriner randevu sequence, siparis sequence, bakici booking activity ve durum diyagramlari.

### Grafik SVG Ciktilari

- `docs/yazilim_muh_final/rendered/charts/`: 5 adet grafik.
- Grafikler: maliyet dagilimi, rol bazli is gucu, faz sureleri, modul kapsam yogunlugu ve proje Gantt plani.

### Ekran Goruntusu Adaylari

- Kok dizindeki `adb_*.png` ve `store_*.png` dosyalari magaza ve mobil dogrulama ekranlari icin adaydir.
- Final Word duzenlemesinde ekran goruntuleri secilerek ilgili modul bolumlerine yerlestirilmelidir.
"""


POSTER_SOURCE = """\
# EvcilHayvan Platformu - Bitirme Projesi Posteri

## Problem

Evcil hayvan sahipleri; pet profili, veteriner randevusu, asi takibi, bakici rezervasyonu, kayip ilanlari, sosyal etkilesim ve magaza ihtiyaclari icin farkli platformlara dagilmaktadir. Bu durum veri tekrarina, guven problemine ve surec takibinde kopukluga neden olmaktadir.

## Cozum

EvcilHayvan Platformu; mobil uygulama, REST API, real-time socket altyapisi, admin paneli ve satici panelinden olusan butunlesik bir dijital ekosistem sunar.

## Ana Moduller

- Kimlik dogrulama ve profil yonetimi
- Evcil hayvan profilleri ve ilanlar
- Veteriner arama, randevu, asi ve saglik gunlugu
- Bakici rezervasyonu, canli konum ve bakim raporu
- Mesajlasma, bildirim ve sosyal akis
- Magaza, sepet, siparis, satici ve admin paneli

## Teknik Mimari

Flutter mobil uygulama, Node.js/Express backend, MongoDB veritabani, Socket.IO gercek zamanli iletisim, Firebase bildirimleri, Google Maps/Places entegrasyonlari ve React tabanli web panelleri kullanilmistir.

## Son Tamamlama

Veteriner filtreleri, randevu yeniden planlama, gercek online gorusme linki, bakici avatar yukleme, bakici finansal ozet, calisma saati kontrolu, duplicate yorum engeli ve fiziksel cihaza APK kurulumu tamamlanmistir.

## Dogrulama

Backend Jest testleri, Flutter widget testleri ve Android debug APK build basariyla tamamlanmistir. Uygulama fiziksel Android cihaza ADB ile kurulmustur.
"""


def clean_markdown(text: str) -> str:
    text = text.replace("\ufeff", "")
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip() + "\n"


def read_source(name: str) -> str:
    path = SOURCE_DIR / name
    return clean_markdown(path.read_text(encoding="utf-8"))


def demote_heading(line: str) -> str:
    if line.startswith("#"):
        hashes, rest = line.split(" ", 1) if " " in line else (line, "")
        level = min(len(hashes) + 1, 6)
        return "#" * level + (" " + rest if rest else "")
    return line


def section_from_file(name: str) -> str:
    content = read_source(name)
    if name == "01_ANA_RAPOR.md":
        lines = content.splitlines()
        start = 0
        end = len(lines)
        for idx, line in enumerate(lines):
            if line.strip() == "## Özet":
                start = idx
                break
        for idx, line in enumerate(lines):
            if line.strip().startswith("## 18. Word"):
                end = idx
                break
        content = "\n".join(lines[start:end]).strip() + "\n"
    title = name.replace(".md", "").replace("_", " ").title()
    lines = [demote_heading(line) for line in content.splitlines()]
    return "\n".join([f"# {title}", "", *lines]).strip() + "\n"


def build_master_markdown() -> str:
    parts = [
        "# T.C. KONYA TEKNIK UNIVERSITESI",
        "## BILGISAYAR MUHENDISLIGI BOLUMU",
        "## BITIRME PROJESI FINAL RAPORU",
        "",
        "| Alan | Bilgi |",
        "|---|---|",
        "| Proje Adi | EvcilHayvan Platformu |",
        "| Ogrenci Adi Soyadi | [AD SOYAD] |",
        "| Ogrenci Numarasi | [NUMARA] |",
        "| Danisman / Ders Yurutucusu | [DANISMAN ADI] |",
        "| Teslim Tarihi | [TARIH] |",
        "",
        "# On Bilgi",
        "",
        "Bu master belge, projede bulunan final rapor paketi, ara rapor kaynaklari, teknik denetim notlari ve son tamamlanan gelistirmeler birlestirilerek hazirlanmistir.",
        "",
        RECENT_COMPLETION.strip(),
        "",
        OLD_REPORT_SUMMARY.strip(),
        "",
        FIGURE_INDEX.strip(),
        "",
    ]
    for name in FINAL_SOURCE_ORDER:
        parts.append(section_from_file(name))
        parts.append("")
    return clean_markdown("\n".join(parts))


def build_poster_markdown() -> str:
    return clean_markdown(POSTER_SOURCE)


def set_cell_shading(cell, fill: str) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:fill"), fill)
    tc_pr.append(shd)


def set_cell_text(cell, text: str, bold: bool = False) -> None:
    cell.text = ""
    p = cell.paragraphs[0]
    p.alignment = WD_ALIGN_PARAGRAPH.LEFT
    run = p.add_run(text.strip())
    run.bold = bold
    run.font.name = "Times New Roman"
    run.font.size = Pt(10)
    cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER


def add_table_from_markdown(doc: Document, rows: list[list[str]]) -> None:
    if not rows:
        return
    width = max(len(r) for r in rows)
    table = doc.add_table(rows=len(rows), cols=width)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.style = "Table Grid"
    for i, row in enumerate(rows):
        for j in range(width):
            cell = table.cell(i, j)
            text = row[j] if j < len(row) else ""
            set_cell_text(cell, text, bold=i == 0)
            if i == 0:
                set_cell_shading(cell, "D9EAD3")
    doc.add_paragraph()


def is_table_separator(line: str) -> bool:
    stripped = line.strip()
    if not stripped.startswith("|") or not stripped.endswith("|"):
        return False
    content = stripped.strip("|").replace(" ", "")
    return bool(content) and all(ch in "-:|" for ch in content)


def parse_table(lines: list[str], start: int) -> tuple[list[list[str]], int]:
    rows: list[list[str]] = []
    i = start
    while i < len(lines):
        line = lines[i].strip()
        if not (line.startswith("|") and line.endswith("|")):
            break
        if not is_table_separator(line):
            rows.append([cell.strip() for cell in line.strip("|").split("|")])
        i += 1
    return rows, i


def strip_inline_markup(text: str) -> str:
    text = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", text)
    text = text.replace("**", "").replace("__", "").replace("`", "")
    return text


def add_markdown_paragraph(doc: Document, text: str) -> None:
    text = strip_inline_markup(text.strip())
    if not text:
        return
    p = doc.add_paragraph(text)
    p.paragraph_format.first_line_indent = Inches(0.25)
    p.paragraph_format.space_after = Pt(6)
    for run in p.runs:
        run.font.name = "Times New Roman"
        run.font.size = Pt(12)


def style_document(doc: Document) -> None:
    section = doc.sections[0]
    section.top_margin = Inches(1)
    section.bottom_margin = Inches(1)
    section.left_margin = Inches(1)
    section.right_margin = Inches(1)
    styles = doc.styles
    styles["Normal"].font.name = "Times New Roman"
    styles["Normal"].font.size = Pt(12)
    for style_name in ["Heading 1", "Heading 2", "Heading 3"]:
        style = styles[style_name]
        style.font.name = "Times New Roman"
        style.font.color.rgb = RGBColor(31, 78, 121)
    styles["Heading 1"].font.size = Pt(16)
    styles["Heading 2"].font.size = Pt(14)
    styles["Heading 3"].font.size = Pt(12)


def markdown_to_docx(md_text: str, out_path: Path) -> None:
    doc = Document()
    style_document(doc)

    lines = md_text.splitlines()
    i = 0
    in_code = False
    code_buffer: list[str] = []

    while i < len(lines):
        line = lines[i]
        stripped = line.strip()
        if stripped.startswith("```"):
            if in_code:
                p = doc.add_paragraph()
                run = p.add_run("\n".join(code_buffer))
                run.font.name = "Consolas"
                run.font.size = Pt(9)
                code_buffer = []
                in_code = False
            else:
                in_code = True
            i += 1
            continue
        if in_code:
            code_buffer.append(line)
            i += 1
            continue
        if not stripped:
            i += 1
            continue
        if stripped.startswith("|") and stripped.endswith("|"):
            rows, next_i = parse_table(lines, i)
            add_table_from_markdown(doc, rows)
            i = next_i
            continue
        if stripped.startswith("# "):
            text = strip_inline_markup(stripped[2:])
            p = doc.add_heading(text, level=1)
            p.alignment = WD_ALIGN_PARAGRAPH.LEFT
            i += 1
            continue
        if stripped.startswith("## "):
            doc.add_heading(strip_inline_markup(stripped[3:]), level=2)
            i += 1
            continue
        if stripped.startswith("### "):
            doc.add_heading(strip_inline_markup(stripped[4:]), level=3)
            i += 1
            continue
        if stripped.startswith("#### "):
            p = doc.add_paragraph()
            run = p.add_run(strip_inline_markup(stripped[5:]))
            run.bold = True
            run.font.name = "Times New Roman"
            run.font.size = Pt(12)
            i += 1
            continue
        if re.match(r"^[-*]\s+", stripped):
            p = doc.add_paragraph(style="List Bullet")
            p.add_run(strip_inline_markup(re.sub(r"^[-*]\s+", "", stripped)))
            i += 1
            continue
        if re.match(r"^\d+\.\s+", stripped):
            p = doc.add_paragraph(style="List Number")
            p.add_run(strip_inline_markup(re.sub(r"^\d+\.\s+", "", stripped)))
            i += 1
            continue
        add_markdown_paragraph(doc, stripped)
        i += 1

    doc.save(out_path)


def add_textbox(slide, x, y, w, h, text, font_size=18, bold=False, color=(30, 30, 30), align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = PptInches(0.08)
    tf.margin_right = PptInches(0.08)
    tf.margin_top = PptInches(0.04)
    tf.margin_bottom = PptInches(0.04)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = PptPt(font_size)
    run.font.bold = bold
    run.font.name = "Arial"
    run.font.color.rgb = PptRGBColor(*color)
    return box


def add_panel(slide, x, y, w, h, title, body, accent=(45, 106, 79), fill=(248, 250, 248)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PptRGBColor(*fill)
    shape.line.color.rgb = PptRGBColor(205, 215, 210)
    shape.line.width = PptPt(1.2)
    add_textbox(slide, x + 0.18, y + 0.15, w - 0.36, 0.35, title, font_size=22, bold=True, color=accent)
    add_textbox(slide, x + 0.18, y + 0.58, w - 0.36, h - 0.75, body, font_size=15, color=(35, 44, 40))
    return shape


def add_module_chip(slide, x, y, w, h, title, subtitle, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PptRGBColor(*fill)
    shape.line.color.rgb = PptRGBColor(255, 255, 255)
    shape.text_frame.clear()
    shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = shape.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title + "\n"
    r.font.name = "Arial"
    r.font.size = PptPt(16)
    r.font.bold = True
    r.font.color.rgb = PptRGBColor(20, 42, 35)
    r2 = p.add_run()
    r2.text = subtitle
    r2.font.name = "Arial"
    r2.font.size = PptPt(11)
    r2.font.color.rgb = PptRGBColor(45, 55, 50)


def create_poster(out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = PptInches(33.11)
    prs.slide_height = PptInches(23.39)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = PptRGBColor(246, 248, 246)

    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(0), PptInches(0), prs.slide_width, PptInches(2.45))
    header.fill.solid()
    header.fill.fore_color.rgb = PptRGBColor(30, 80, 63)
    header.line.color.rgb = PptRGBColor(30, 80, 63)
    add_textbox(slide, 0.85, 0.35, 21.5, 0.75, "EvcilHayvan Platformu", 34, True, (255, 255, 255))
    add_textbox(slide, 0.88, 1.12, 23.5, 0.55, "Bitirme Projesi Final Raporu ve Akademik Poster Ozeti", 20, False, (225, 240, 232))
    add_textbox(slide, 24.2, 0.52, 7.8, 1.0, "Ogrenci: [AD SOYAD]\nDanisman: [DANISMAN ADI]\nTarih: [TARIH]", 14, False, (255, 255, 255), PP_ALIGN.RIGHT)

    add_panel(
        slide, 0.85, 3.05, 7.7, 4.0, "Problem",
        "Evcil hayvan sahipleri veteriner, bakici, magaza, kayip ilanlari, sosyal etkilesim ve saglik takibi icin farkli uygulamalara dagilmaktadir. Bu daginik yapi veri tekrarina, takip zorluguna ve guven sorununa neden olur.",
        accent=(33, 93, 72), fill=(255, 255, 255),
    )
    add_panel(
        slide, 0.85, 7.45, 7.7, 4.35, "Cozum Yaklasimi",
        "Proje; mobil uygulama, REST API, gercek zamanli socket katmani, MongoDB veritabani, admin paneli ve satici panelini tek platformda birlestiren entegre bir ekosistem olarak tasarlanmistir.",
        accent=(33, 93, 72), fill=(255, 255, 255),
    )
    add_panel(
        slide, 0.85, 12.25, 7.7, 4.3, "Teknoloji Yigini",
        "Mobil: Flutter\nBackend: Node.js / Express\nVeritabani: MongoDB\nReal-time: Socket.IO\nBildirim: Firebase\nHarita: Google Maps / Places\nPaneller: React",
        accent=(33, 93, 72), fill=(255, 255, 255),
    )

    add_panel(
        slide, 9.15, 3.05, 14.65, 2.2, "Sistem Mimarisi",
        "Flutter mobil istemci ve web panelleri, Express REST API uzerinden servis katmanina baglanir. MongoDB kalici veri, Socket.IO anlik konum/mesaj/bildirim akislarini, Firebase ise push bildirimlerini destekler.",
        accent=(121, 81, 32), fill=(255, 252, 246),
    )
    add_textbox(slide, 9.25, 5.72, 14.2, 0.45, "Ana Moduller", 24, True, (30, 80, 63), PP_ALIGN.CENTER)
    colors = [(219, 235, 226), (232, 241, 255), (255, 239, 218), (235, 229, 248), (232, 244, 244), (246, 235, 225)]
    modules = [
        ("Pet Profili", "Kimlik, tur, saglik, ilan"),
        ("Veteriner", "Arama, filtre, randevu, asi"),
        ("Bakici", "Rezervasyon, konum, rapor"),
        ("Magaza", "Urun, sepet, siparis"),
        ("Mesaj", "Sohbet ve bildirim"),
        ("Admin/Satici", "Panel, denetim, finans"),
    ]
    for idx, (title, subtitle) in enumerate(modules):
        row, col = divmod(idx, 3)
        add_module_chip(slide, 9.45 + col * 4.65, 6.45 + row * 2.05, 4.15, 1.55, title, subtitle, colors[idx])

    # Simple native proof bars.
    add_panel(slide, 9.15, 11.0, 14.65, 5.55, "Kapsam ve Dogrulama", "", accent=(33, 93, 72), fill=(255, 255, 255))
    add_textbox(slide, 9.55, 11.72, 5.5, 0.45, "Tamamlanan Alanlar", 17, True, (33, 93, 72))
    proof = [("Backend test", 10), ("Flutter widget", 3), ("Diyagram", 25), ("Grafik", 5), ("APK kurulum", 1)]
    max_val = 25
    for i, (label, value) in enumerate(proof):
        y = 12.35 + i * 0.65
        add_textbox(slide, 9.55, y, 3.2, 0.32, label, 12, False, (45, 55, 50))
        slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(12.95), PptInches(y + 0.06), PptInches(5.8), PptInches(0.22)).fill.solid()
        bg_bar = slide.shapes[-1]
        bg_bar.fill.fore_color.rgb = PptRGBColor(226, 232, 226)
        bg_bar.line.color.rgb = PptRGBColor(226, 232, 226)
        fg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(12.95), PptInches(y + 0.06), PptInches(5.8 * value / max_val), PptInches(0.22))
        fg.fill.solid()
        fg.fill.fore_color.rgb = PptRGBColor(45, 106, 79)
        fg.line.color.rgb = PptRGBColor(45, 106, 79)
        add_textbox(slide, 19.0, y - 0.02, 1.0, 0.32, str(value), 12, True, (30, 80, 63))
    add_textbox(slide, 9.55, 15.3, 13.5, 0.7, "Son dogrulama: Backend Jest testleri, Flutter widget testleri, debug APK build ve fiziksel cihaza ADB kurulumu basarili.", 13, False, (40, 50, 45))

    add_panel(
        slide, 24.4, 3.05, 7.8, 4.35, "Son Gelistirmeler",
        "Veteriner filtreleri, randevu yeniden planlama, gercek online gorusme linki, bakici avatar yukleme, bakici finansal ozet, calisma saati kontrolu ve duplicate yorum engeli tamamlandi.",
        accent=(121, 81, 32), fill=(255, 252, 246),
    )
    add_panel(
        slide, 24.4, 7.85, 7.8, 4.1, "Proje Katkisi",
        "Tekil evcil hayvan uygulamalarinin parcalı yapi sorununu azaltir. Kullanici, veteriner, bakici, satici ve yonetici rollerini ayni ekosistemde bulusturur.",
        accent=(121, 81, 32), fill=(255, 252, 246),
    )
    add_panel(
        slide, 24.4, 12.4, 7.8, 4.15, "Gelecek Calismalar",
        "Google Meet/Calendar OAuth entegrasyonu, odeme altyapisi, daha kapsamli saha testleri, raporlama paneli ve performans izleme sureclerinin genisletilmesi hedeflenmektedir.",
        accent=(121, 81, 32), fill=(255, 252, 246),
    )

    footer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(0), PptInches(21.75), prs.slide_width, PptInches(1.64))
    footer.fill.solid()
    footer.fill.fore_color.rgb = PptRGBColor(235, 241, 237)
    footer.line.color.rgb = PptRGBColor(235, 241, 237)
    add_textbox(slide, 0.85, 22.08, 31.4, 0.52, "Cikti paketleri: BITIRME_FINAL_MASTER.md | Bitirme_Final_Raporu.docx | BITIRME_FINAL_POSTER.md | Bitirme_Final_Posteri.pptx", 14, False, (30, 80, 63), PP_ALIGN.CENTER)

    prs.save(out_path)


def write_notes() -> None:
    text = """# Teslim Notlari

Bu klasor, tum mevcut rapor kaynaklarindan uretilen final teslim taslaklarini icerir.

## Uretilen Dosyalar

- `BITIRME_FINAL_MASTER.md`: Tek master final rapor kaynagi.
- `Bitirme_Final_Raporu.docx`: Word final rapor taslagi.
- `BITIRME_FINAL_POSTER.md`: Poster metin kaynagi.
- `Bitirme_Final_Posteri.pptx`: Tek sayfalik akademik poster taslagi.

## Doldurulmasi Gereken Alanlar

- Ogrenci adi soyadi
- Ogrenci numarasi
- Danisman / ders yurutucusu
- Teslim tarihi
- Gerekirse universite/fakulte bilgisi

## Teknik Not

Bu makinede Pandoc, LibreOffice/soffice ve artifact-tool bulunmadigi icin DOCX/PPTX yapisal olarak uretilmistir; LibreOffice ile sayfa sayfa render QA yapilamamistir. Final teslimden once Word veya PowerPoint icinde son gorsel kontrol onerilir.
"""
    NOTES_MD.write_text(text, encoding="utf-8")


def verify_zip(path: Path) -> None:
    with zipfile.ZipFile(path) as zf:
        bad = zf.testzip()
        if bad:
            raise RuntimeError(f"Corrupt file entry in {path}: {bad}")


def main() -> None:
    FINAL_DIR.mkdir(parents=True, exist_ok=True)
    master = build_master_markdown()
    poster = build_poster_markdown()
    MASTER_MD.write_text(master, encoding="utf-8")
    POSTER_MD.write_text(poster, encoding="utf-8")
    markdown_to_docx(master, REPORT_DOCX)
    create_poster(POSTER_PPTX)
    write_notes()
    verify_zip(REPORT_DOCX)
    verify_zip(POSTER_PPTX)
    print(f"Wrote {MASTER_MD}")
    print(f"Wrote {REPORT_DOCX}")
    print(f"Wrote {POSTER_MD}")
    print(f"Wrote {POSTER_PPTX}")
    print(f"Wrote {NOTES_MD}")


if __name__ == "__main__":
    main()
